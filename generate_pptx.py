from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree
import copy

# ─── COLOR PALETTE ────────────────────────────────────────────────────────────
PINK        = RGBColor(0xE0, 0x1A, 0x8C)   # #E01A8C  – magenta principal
PINK_DARK   = RGBColor(0x9C, 0x0F, 0x62)   # #9C0F62  – magenta escuro
LIME        = RGBColor(0x7E, 0xD3, 0x21)   # #7ED321  – verde limão
LIME_DARK   = RGBColor(0x55, 0x91, 0x12)   # #559112  – verde escuro
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF5, 0xF5, 0xF5)
DARK        = RGBColor(0x12, 0x12, 0x20)   # fundo escuro
DARK2       = RGBColor(0x1E, 0x1E, 0x32)   # card escuro
GRAY_LIGHT  = RGBColor(0xCC, 0xCC, 0xCC)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─── HELPERS ──────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb=None, alpha=None, line_rgb=None, line_width=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = line_width or Pt(1)
    else:
        line.fill.background()
    return shape


def add_text_box(slide, text, x, y, w, h,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_label_value(slide, label, value, x, y, label_color=LIME, value_color=WHITE):
    """Two-line block: label (lime small) + value (white bold)."""
    add_text_box(slide, label, x, y, Inches(3.5), Inches(0.35),
                 font_size=10, bold=False, color=label_color, align=PP_ALIGN.LEFT)
    add_text_box(slide, value, x, y + Inches(0.32), Inches(3.5), Inches(0.45),
                 font_size=16, bold=True, color=value_color, align=PP_ALIGN.LEFT)


def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def diagonal_bg(slide, color_left=PINK, color_right=LIME):
    """Simula fundo diagonal com dois triângulos via freeform."""
    set_bg(slide, DARK)
    # Retângulo esquerdo grande (pink)
    add_rect(slide, Inches(0), Inches(0), Inches(7.5), SLIDE_H, fill_rgb=PINK_DARK)
    # Retângulo direito (lime)
    add_rect(slide, Inches(7.5), Inches(0), Inches(5.83), SLIDE_H, fill_rgb=LIME_DARK)
    # Faixa diagonal (override com shape inclinado — usamos paralelogramo via freeform)
    # Simplificado: faixa no meio para dar efeito de transição
    add_rect(slide, Inches(6.8), Inches(0), Inches(1.4), SLIDE_H, fill_rgb=PINK)
    add_rect(slide, Inches(7.2), Inches(0), Inches(1.0), SLIDE_H, fill_rgb=LIME)


def pill_tag(slide, text, x, y, bg=LIME, fg=DARK, size=10):
    """Pequena tag arredondada (usando retângulo por limitação da lib)."""
    w = Inches(len(text) * 0.095 + 0.3)
    r = add_rect(slide, x, y, w, Inches(0.32), fill_rgb=bg)
    tf = r.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg
    run.font.name = "Calibri"
    return r


def section_header_bar(slide, title):
    """Faixa de cabeçalho padrão para slides internos."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.0), fill_rgb=DARK2)
    add_rect(slide, 0, 0, Inches(0.12), Inches(1.0), fill_rgb=PINK)
    add_text_box(slide, title, Inches(0.3), Inches(0.15), Inches(12), Inches(0.7),
                 font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)


def accent_line(slide, x, y, w, color=PINK):
    add_rect(slide, x, y, w, Inches(0.045), fill_rgb=color)


# ─── SLIDES ───────────────────────────────────────────────────────────────────

def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Fundo: pink escuro à esquerda, lime escuro à direita
    set_bg(slide, DARK)
    add_rect(slide, 0, 0, Inches(8.0), SLIDE_H, fill_rgb=PINK_DARK)
    add_rect(slide, Inches(7.5), 0, Inches(5.83), SLIDE_H, fill_rgb=LIME_DARK)
    # Faixa diagonal central
    add_rect(slide, Inches(7.0), 0, Inches(1.5), SLIDE_H, fill_rgb=PINK)
    add_rect(slide, Inches(7.4), 0, Inches(1.0), SLIDE_H, fill_rgb=LIME)

    # Decorativos: "hashtag" círculos
    for pos in [(Inches(0.3), Inches(0.3)), (Inches(1.2), Inches(5.8)),
                (Inches(5.5), Inches(6.8)), (Inches(11.5), Inches(0.5))]:
        add_text_box(slide, "#", pos[0], pos[1], Inches(0.8), Inches(0.8),
                     font_size=32, bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
        # Pequeno círculo decorativo como fundo
        c = slide.shapes.add_shape(9, pos[0]-Inches(0.05), pos[1]-Inches(0.05),
                                   Inches(0.7), Inches(0.7))  # 9 = oval
        c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF)
        c.line.fill.background()
        sp = c._element; sp.getparent().remove(sp)  # remove círculo, manter só texto

    # Estrelas decorativas (pontos)
    for pos, sz in [((Inches(0.8), Inches(1.5)), 8), ((Inches(4.0), Inches(0.4)), 6),
                    ((Inches(6.2), Inches(3.0)), 10), ((Inches(10.0), Inches(6.5)), 7)]:
        add_text_box(slide, "✦", pos[0], pos[1], Inches(0.5), Inches(0.5),
                     font_size=sz*2, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

    # Tag projeto
    pill_tag(slide, "PROJETO PESSOAL", Inches(1.0), Inches(1.8), bg=LIME, fg=DARK, size=11)

    # Título principal
    add_text_box(slide, "Local", Inches(1.0), Inches(2.35), Inches(6), Inches(1.1),
                 font_size=64, bold=True, color=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri")
    add_text_box(slide, "Data Stack", Inches(1.0), Inches(3.1), Inches(6), Inches(1.2),
                 font_size=64, bold=True, color=LIME, align=PP_ALIGN.LEFT, font_name="Calibri")

    # Subtítulo
    add_text_box(slide,
                 "Plataforma enterprise de dados rodando 100% local.\nCusto zero. Ambiente real.",
                 Inches(1.0), Inches(4.4), Inches(5.5), Inches(1.0),
                 font_size=16, bold=False, color=OFF_WHITE, align=PP_ALIGN.LEFT)

    # Linha decorativa
    accent_line(slide, Inches(1.0), Inches(4.25), Inches(3.5), color=LIME)

    # Lado direito: 3 métricas rápidas
    for i, (label, val) in enumerate([
        ("FERRAMENTAS", "14+"),
        ("FASES",       "7"),
        ("CUSTO/MÊS",   "R$ 0"),
    ]):
        bx = Inches(9.5)
        by = Inches(1.8 + i * 1.6)
        add_rect(slide, bx, by, Inches(2.5), Inches(1.3), fill_rgb=DARK2)
        accent_line(slide, bx, by, Inches(2.5), color=LIME if i % 2 == 0 else PINK)
        add_text_box(slide, val, bx + Inches(0.2), by + Inches(0.1), Inches(2.1), Inches(0.75),
                     font_size=40, bold=True, color=LIME if i % 2 == 0 else PINK, align=PP_ALIGN.LEFT)
        add_text_box(slide, label, bx + Inches(0.2), by + Inches(0.85), Inches(2.1), Inches(0.35),
                     font_size=10, bold=False, color=GRAY_LIGHT, align=PP_ALIGN.LEFT)

    return slide


def slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    section_header_bar(slide, "Visão Geral")
    accent_line(slide, 0, Inches(1.0), SLIDE_W, color=PINK)

    # Bloco objetivo
    add_rect(slide, Inches(0.4), Inches(1.2), Inches(5.8), Inches(2.0), fill_rgb=DARK2)
    accent_line(slide, Inches(0.4), Inches(1.2), Inches(5.8), color=PINK)
    add_text_box(slide, "OBJETIVO", Inches(0.65), Inches(1.3), Inches(5.4), Inches(0.4),
                 font_size=11, bold=True, color=PINK)
    add_text_box(slide,
                 "Consolidar skills de engenheiro de dados sênior simulando um ambiente de produção enterprise real — sem gastar nada.",
                 Inches(0.65), Inches(1.65), Inches(5.3), Inches(0.9),
                 font_size=14, bold=False, color=WHITE)

    # Bloco motivação
    add_rect(slide, Inches(7.0), Inches(1.2), Inches(5.8), Inches(2.0), fill_rgb=DARK2)
    accent_line(slide, Inches(7.0), Inches(1.2), Inches(5.8), color=LIME)
    add_text_box(slide, "MOTIVAÇÃO", Inches(7.25), Inches(1.3), Inches(5.4), Inches(0.4),
                 font_size=11, bold=True, color=LIME)
    add_text_box(slide,
                 "Hardware potente disponível localmente (i9 14ª gen, 32 GB RAM, RTX 5070 Ti 16 GB VRAM) — por que não usá-lo como data center pessoal?",
                 Inches(7.25), Inches(1.65), Inches(5.3), Inches(0.9),
                 font_size=14, bold=False, color=WHITE)

    # Três princípios
    principles = [
        (PINK,  "CUSTO ZERO",     "Somente ferramentas open-source.\nSem cloud, sem licença."),
        (LIME,  "100% LOCAL",     "Docker Compose com profiles.\nTudo na sua máquina."),
        (WHITE, "ENTERPRISE-LIKE","Simula pipelines reais:\nbatch, streaming e IA."),
    ]
    for i, (color, title, desc) in enumerate(principles):
        bx = Inches(0.4 + i * 4.3)
        by = Inches(3.6)
        add_rect(slide, bx, by, Inches(3.9), Inches(2.8), fill_rgb=DARK2)
        accent_line(slide, bx, by, Inches(3.9), color=color)
        add_text_box(slide, title, bx + Inches(0.25), by + Inches(0.2), Inches(3.4), Inches(0.5),
                     font_size=15, bold=True, color=color)
        add_text_box(slide, desc, bx + Inches(0.25), by + Inches(0.7), Inches(3.4), Inches(1.8),
                     font_size=13, bold=False, color=OFF_WHITE)

    return slide


def slide_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    section_header_bar(slide, "Stack Tecnológica")
    accent_line(slide, 0, Inches(1.0), SLIDE_W, color=LIME)

    categories = [
        ("Orquestração",   ["Airflow 3+"],                             PINK),
        ("Storage",        ["MinIO", "Apache Iceberg"],                LIME),
        ("Lakehouse",      ["Dremio Community", "Project Nessie"],     PINK),
        ("DWH",            ["ClickHouse"],                             LIME),
        ("Transformação",  ["dbt Core"],                               PINK),
        ("Qualidade",      ["Great Expectations"],                     LIME),
        ("Streaming",      ["Redpanda (Kafka-compat.)"],               PINK),
        ("Ingestão",       ["Airbyte OSS", "Python custom"],           LIME),
        ("Catálogo",       ["OpenMetadata"],                           PINK),
        ("Visualização",   ["Apache Superset"],                        LIME),
        ("Observabilidade",["Grafana", "Prometheus"],                  PINK),
        ("Agente IA",      ["Ollama", "FastAPI", "LangChain","ChromaDB"], LIME),
    ]

    cols = 4
    rows = 3
    cw = Inches(3.1)
    ch = Inches(1.7)
    margin_x = Inches(0.25)
    margin_y = Inches(1.2)
    gap_x = Inches(0.1)
    gap_y = Inches(0.12)

    for idx, (cat, tools, color) in enumerate(categories):
        col = idx % cols
        row = idx // cols
        bx = margin_x + col * (cw + gap_x)
        by = margin_y + row * (ch + gap_y)
        add_rect(slide, bx, by, cw, ch, fill_rgb=DARK2)
        accent_line(slide, bx, by, cw, color=color)
        add_text_box(slide, cat.upper(), bx + Inches(0.15), by + Inches(0.05),
                     cw - Inches(0.3), Inches(0.35),
                     font_size=9, bold=True, color=color)
        for t_i, tool in enumerate(tools):
            add_text_box(slide, f"• {tool}", bx + Inches(0.15),
                         by + Inches(0.38 + t_i * 0.32),
                         cw - Inches(0.3), Inches(0.35),
                         font_size=11, bold=False, color=WHITE)

    return slide


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    section_header_bar(slide, "Arquitetura Medallion")
    accent_line(slide, 0, Inches(1.0), SLIDE_W, color=PINK)

    # Diagrama de fluxo: Fonte → Bronze → Silver → Gold → Consumo
    stages = [
        ("FONTES",   "APIs, DBs,\nArquivos, Streaming",  DARK2,  GRAY_LIGHT),
        ("BRONZE",   "Raw data\nIngested as-is",          RGBColor(0x8B,0x45,0x13), RGBColor(0xFF,0xA5,0x00)),
        ("SILVER",   "Dados limpos\ne padronizados",      RGBColor(0x1C,0x3A,0x52), RGBColor(0x87,0xCE,0xFA)),
        ("GOLD",     "Dados prontos\npara consumo",       RGBColor(0x3A,0x2F,0x05), RGBColor(0xFF,0xD7,0x00)),
        ("CONSUMO",  "BI, APIs,\nIA Agents",              DARK2,  LIME),
    ]

    n = len(stages)
    bw = Inches(2.2)
    bh = Inches(3.0)
    gap = Inches(0.28)
    total_w = n * bw + (n - 1) * gap
    start_x = (SLIDE_W - total_w) / 2
    by = Inches(1.6)

    for i, (name, desc, bg, accent) in enumerate(stages):
        bx = start_x + i * (bw + gap)
        add_rect(slide, bx, by, bw, bh, fill_rgb=bg)
        accent_line(slide, bx, by, bw, color=accent)
        add_text_box(slide, name, bx + Inches(0.1), by + Inches(0.1), bw - Inches(0.2), Inches(0.5),
                     font_size=14, bold=True, color=accent, align=PP_ALIGN.CENTER)
        add_text_box(slide, desc, bx + Inches(0.1), by + Inches(0.65), bw - Inches(0.2), Inches(1.8),
                     font_size=11, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

        # Seta entre blocos
        if i < n - 1:
            ax = bx + bw + Inches(0.02)
            ay = by + bh / 2 - Inches(0.15)
            add_text_box(slide, "→", ax, ay, gap, Inches(0.4),
                         font_size=20, bold=True, color=LIME, align=PP_ALIGN.CENTER)

    # Ferramentas associadas por camada
    tools_row = [
        ("", ""),
        ("Airbyte OSS\nRedpanda", "Ingestão"),
        ("dbt Core\nGreat Expectations", "Transformação"),
        ("Dremio + Nessie\nClickHouse", "Lakehouse/DWH"),
        ("Superset\nAI Agent", ""),
    ]
    for i, (tools, label) in enumerate(tools_row):
        if not tools:
            continue
        bx = start_x + i * (bw + gap)
        add_text_box(slide, tools, bx + Inches(0.1), by + bh + Inches(0.15),
                     bw - Inches(0.2), Inches(0.9),
                     font_size=9, bold=False, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)

    # Legenda de storage
    add_rect(slide, Inches(0.4), Inches(5.7), Inches(12.5), Inches(0.7), fill_rgb=DARK2)
    add_text_box(slide, "Storage transversal:",
                 Inches(0.6), Inches(5.75), Inches(2.2), Inches(0.5),
                 font_size=11, bold=True, color=PINK)
    add_text_box(slide, "MinIO (object storage S3-compatible)  +  Apache Iceberg (table format)  →  gerenciado pelo Project Nessie (data catalog)",
                 Inches(2.8), Inches(5.75), Inches(10), Inches(0.5),
                 font_size=11, bold=False, color=OFF_WHITE)

    return slide


def slide_phases(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    section_header_bar(slide, "Roadmap — 7 Fases")
    accent_line(slide, 0, Inches(1.0), SLIDE_W, color=LIME)

    phases = [
        ("1", "Fundação",             "Docker Compose profiles\nCore services + 1ª ingestão",  PINK),
        ("2", "Lakehouse",            "MinIO + Iceberg + Nessie\nDremio Community",             LIME),
        ("3", "Catálogo & Governança","OpenMetadata\nLinhagem e metadados",                    PINK),
        ("4", "Streaming",            "Redpanda (Kafka)\nPipelines em tempo real",              LIME),
        ("5", "Viz & Observabilidade","Superset + Grafana\nPrometheus dashboards",             PINK),
        ("6", "Agente IA",            "Ollama + LangChain\nText2SQL & RAG",                    LIME),
        ("7", "CI/CD & Maturidade",   "GitHub Actions\nTestes, lint, qualidade",               PINK),
    ]

    cols = 4
    bw = Inches(3.0)
    bh = Inches(2.2)
    gap_x = Inches(0.18)
    gap_y = Inches(0.18)
    start_x = Inches(0.35)
    start_y = Inches(1.25)

    for i, (num, title, desc, color) in enumerate(phases):
        col = i % cols
        row = i // cols
        bx = start_x + col * (bw + gap_x)
        by = start_y + row * (bh + gap_y)
        add_rect(slide, bx, by, bw, bh, fill_rgb=DARK2)
        accent_line(slide, bx, by, bw, color=color)
        # Número grande
        add_text_box(slide, num, bx + Inches(0.12), by + Inches(0.05),
                     Inches(0.6), Inches(0.8),
                     font_size=36, bold=True, color=color, align=PP_ALIGN.LEFT)
        # Título fase
        add_text_box(slide, title, bx + Inches(0.65), by + Inches(0.1),
                     bw - Inches(0.75), Inches(0.6),
                     font_size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        # Descrição
        add_text_box(slide, desc, bx + Inches(0.15), by + Inches(0.85),
                     bw - Inches(0.3), Inches(1.2),
                     font_size=11, bold=False, color=GRAY_LIGHT, align=PP_ALIGN.LEFT)

    return slide


def slide_ai_agent(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    section_header_bar(slide, "Agente IA — LLM Local")
    accent_line(slide, 0, Inches(1.0), SLIDE_W, color=LIME)

    # Badge modelo
    add_rect(slide, Inches(0.4), Inches(1.1), Inches(5.5), Inches(0.5), fill_rgb=DARK2)
    add_text_box(slide, "Modelo: Ollama · Qwen2.5-Coder 14B Q4_K_M  |  Backend: FastAPI · LangChain · ChromaDB",
                 Inches(0.55), Inches(1.15), Inches(5.2), Inches(0.4),
                 font_size=10, bold=False, color=GRAY_LIGHT)

    # Caso 1
    bx1, by1 = Inches(0.4), Inches(1.85)
    bw1, bh1 = Inches(5.9), Inches(4.8)
    add_rect(slide, bx1, by1, bw1, bh1, fill_rgb=DARK2)
    accent_line(slide, bx1, by1, bw1, color=PINK)
    pill_tag(slide, "CASO DE USO 1", bx1 + Inches(0.2), by1 + Inches(0.15), bg=PINK, fg=WHITE, size=9)
    add_text_box(slide, "Text2SQL Conversacional", bx1 + Inches(0.2), by1 + Inches(0.55),
                 bw1 - Inches(0.4), Inches(0.55), font_size=18, bold=True, color=WHITE)
    items1 = [
        "Usuário faz perguntas em linguagem natural",
        "Agente consulta metadados do OpenMetadata",
        "Gera SQL otimizado para ClickHouse",
        "Retorna resposta com contexto de negócio",
    ]
    for j, item in enumerate(items1):
        add_text_box(slide, f"→  {item}", bx1 + Inches(0.2), by1 + Inches(1.3 + j * 0.65),
                     bw1 - Inches(0.4), Inches(0.55), font_size=12, color=OFF_WHITE)

    # Caso 2
    bx2, by2 = Inches(7.0), Inches(1.85)
    bw2, bh2 = Inches(5.9), Inches(4.8)
    add_rect(slide, bx2, by2, bw2, bh2, fill_rgb=DARK2)
    accent_line(slide, bx2, by2, bw2, color=LIME)
    pill_tag(slide, "CASO DE USO 2", bx2 + Inches(0.2), by2 + Inches(0.15), bg=LIME, fg=DARK, size=9)
    add_text_box(slide, "Architecture Assistant (RAG)", bx2 + Inches(0.2), by2 + Inches(0.55),
                 bw2 - Inches(0.4), Inches(0.55), font_size=18, bold=True, color=WHITE)
    items2 = [
        "Base de conhecimento: PRD + ADRs internos",
        "ChromaDB como vector store local",
        "Responde dúvidas sobre decisões de arquitetura",
        "Sem dados saindo da máquina local",
    ]
    for j, item in enumerate(items2):
        add_text_box(slide, f"→  {item}", bx2 + Inches(0.2), by2 + Inches(1.3 + j * 0.65),
                     bw2 - Inches(0.4), Inches(0.55), font_size=12, color=OFF_WHITE)

    return slide


def slide_infra(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    section_header_bar(slide, "Infraestrutura Local")
    accent_line(slide, 0, Inches(1.0), SLIDE_W, color=PINK)

    # Hardware specs
    hw_specs = [
        ("CPU",    "Intel Core i9  14ª Geração",    PINK),
        ("RAM",    "32 GB",                          LIME),
        ("GPU",    "RTX 5070 Ti — 16 GB VRAM",       PINK),
        ("OS",     "WSL2 / Linux",                   LIME),
    ]
    add_text_box(slide, "HARDWARE", Inches(0.4), Inches(1.2), Inches(5.5), Inches(0.4),
                 font_size=11, bold=True, color=GRAY_LIGHT)
    for i, (key, val, color) in enumerate(hw_specs):
        by = Inches(1.65 + i * 0.85)
        add_rect(slide, Inches(0.4), by, Inches(5.5), Inches(0.72), fill_rgb=DARK2)
        accent_line(slide, Inches(0.4), by, Inches(0.12), color=color)
        add_text_box(slide, key, Inches(0.65), by + Inches(0.08), Inches(1.5), Inches(0.35),
                     font_size=10, bold=True, color=color)
        add_text_box(slide, val, Inches(2.1), by + Inches(0.08), Inches(3.7), Inches(0.55),
                     font_size=14, bold=True, color=WHITE)

    # Docker profiles
    add_text_box(slide, "DOCKER COMPOSE PROFILES", Inches(6.8), Inches(1.2), Inches(6.0), Inches(0.4),
                 font_size=11, bold=True, color=GRAY_LIGHT)

    profiles = [
        ("core",      ["Airflow", "MinIO", "ClickHouse", "Redpanda", "Airbyte"], PINK,
         "Fundação — sempre ativo"),
        ("analytics", ["Dremio", "Nessie", "dbt", "Great Expectations",
                       "OpenMetadata", "Superset"],                               LIME,
         "Camada analítica completa"),
        ("ai",        ["Ollama", "FastAPI", "LangChain", "ChromaDB"],            RGBColor(0xBB,0x86,0xFC),
         "Stack de IA generativa"),
    ]
    for i, (name, svcs, color, label) in enumerate(profiles):
        by = Inches(1.65 + i * 1.85)
        add_rect(slide, Inches(6.8), by, Inches(6.0), Inches(1.65), fill_rgb=DARK2)
        accent_line(slide, Inches(6.8), by, Inches(6.0), color=color)
        add_text_box(slide, f"--profile {name}", Inches(7.0), by + Inches(0.05),
                     Inches(2.5), Inches(0.4), font_size=12, bold=True, color=color)
        add_text_box(slide, label, Inches(9.5), by + Inches(0.05), Inches(3.1), Inches(0.4),
                     font_size=10, bold=False, color=GRAY_LIGHT, align=PP_ALIGN.RIGHT)
        svc_text = "  ·  ".join(svcs)
        add_text_box(slide, svc_text, Inches(7.0), by + Inches(0.5), Inches(5.6), Inches(0.9),
                     font_size=11, bold=False, color=OFF_WHITE)

    return slide


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_rect(slide, 0, 0, Inches(6.5), SLIDE_H, fill_rgb=PINK_DARK)
    add_rect(slide, Inches(6.0), 0, Inches(7.33), SLIDE_H, fill_rgb=LIME_DARK)
    add_rect(slide, Inches(5.8), 0, Inches(1.2), SLIDE_H, fill_rgb=PINK)
    add_rect(slide, Inches(6.2), 0, Inches(0.8), SLIDE_H, fill_rgb=LIME)

    add_text_box(slide, "Build it.", Inches(0.8), Inches(2.0), Inches(5.2), Inches(1.1),
                 font_size=60, bold=True, color=WHITE, font_name="Calibri")
    add_text_box(slide, "Learn it.", Inches(0.8), Inches(2.9), Inches(5.2), Inches(1.1),
                 font_size=60, bold=True, color=LIME, font_name="Calibri")
    add_text_box(slide, "Own it.", Inches(0.8), Inches(3.8), Inches(5.2), Inches(1.1),
                 font_size=60, bold=True, color=WHITE, font_name="Calibri")

    accent_line(slide, Inches(0.8), Inches(5.1), Inches(4.0), color=LIME)

    add_text_box(slide, "Local Data Stack  —  Projeto Pessoal",
                 Inches(0.8), Inches(5.3), Inches(5.0), Inches(0.5),
                 font_size=13, bold=False, color=GRAY_LIGHT)

    # Direita: resumo tags
    tags = ["Airflow", "MinIO", "Iceberg", "Dremio", "ClickHouse",
            "dbt", "Redpanda", "Airbyte", "OpenMetadata", "Superset",
            "Grafana", "Ollama", "LangChain", "ChromaDB"]
    colors_cycle = [PINK, LIME, WHITE]
    tx, ty = Inches(7.5), Inches(1.5)
    for j, tag in enumerate(tags):
        col = j % 3
        row = j // 3
        pill_tag(slide, tag,
                 tx + col * Inches(1.9),
                 ty + row * Inches(0.55),
                 bg=colors_cycle[col], fg=DARK if colors_cycle[col] != WHITE else DARK, size=10)

    return slide


# ─── MAIN ─────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_overview(prs)
    slide_stack(prs)
    slide_architecture(prs)
    slide_phases(prs)
    slide_ai_agent(prs)
    slide_infra(prs)
    slide_closing(prs)

    out = "samples/local_data_stack.pptx"
    prs.save(out)
    print(f"✓ Saved → {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
